#!/usr/bin/env python3
"""
Análisis de Engagement Avanzado - Funcionalidades Premium
=========================================================
Mejoras adicionales premium:
- Exportación a PowerPoint con visualizaciones
- Dashboard web interactivo
- Análisis de competencia
- Sistema de alertas automáticas
- Integración con herramientas externas
- Análisis de ROI avanzado
- Generación automática de reportes ejecutivos
"""

import os
import sys
import json
import logging
from typing import Dict, List, Any, Optional
from datetime import datetime, timedelta
from collections import defaultdict
import statistics

try:
    from analisis_engagement_contenido import AnalizadorEngagement, Publicacion
    from analisis_engagement_ai import AnalizadorEngagementAI
    from analisis_engagement_mejorado import AnalizadorEngagementMejorado
except ImportError:
    print("Error: Módulos de análisis no encontrados")
    sys.exit(1)

logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)


class AnalizadorEngagementAvanzado:
    """Analizador de engagement con funcionalidades premium avanzadas"""
    
    def __init__(
        self,
        analizador_base: AnalizadorEngagement,
        analizador_ai: Optional[AnalizadorEngagementAI] = None,
        analizador_mejorado: Optional[AnalizadorEngagementMejorado] = None
    ):
        """
        Inicializa el analizador avanzado
        
        Args:
            analizador_base: Instancia del AnalizadorEngagement base
            analizador_ai: Instancia opcional del AnalizadorEngagementAI
            analizador_mejorado: Instancia opcional del AnalizadorEngagementMejorado
        """
        self.analizador = analizador_base
        self.analizador_ai = analizador_ai or AnalizadorEngagementAI(analizador_base)
        self.analizador_mejorado = analizador_mejorado or AnalizadorEngagementMejorado(analizador_base, self.analizador_ai)
    
    def exportar_powerpoint(self, reporte: Dict[str, Any], output_file: str) -> Dict[str, Any]:
        """
        Exporta el análisis a PowerPoint con visualizaciones
        
        Args:
            reporte: Reporte generado
            output_file: Ruta del archivo PowerPoint a generar
        
        Returns:
            Información del archivo generado
        """
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            from pptx.enum.text import PP_ALIGN
            from pptx.dml.color import RGBColor
        except ImportError:
            return {
                "error": "python-pptx no está instalado. Instálalo con: pip install python-pptx"
            }
        
        try:
            prs = Presentation()
            prs.slide_width = Inches(10)
            prs.slide_height = Inches(7.5)
            
            # Slide 1: Portada
            slide_portada = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
            title_shape = slide_portada.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1.5))
            title_frame = title_shape.text_frame
            title_frame.text = "Análisis de Engagement"
            title_frame.paragraphs[0].font.size = Pt(44)
            title_frame.paragraphs[0].font.bold = True
            
            subtitle_shape = slide_portada.shapes.add_textbox(Inches(1), Inches(4), Inches(8), Inches(1))
            subtitle_frame = subtitle_shape.text_frame
            subtitle_frame.text = f"Reporte generado: {datetime.now().strftime('%d/%m/%Y %H:%M')}"
            subtitle_frame.paragraphs[0].font.size = Pt(18)
            
            # Slide 2: Resumen Ejecutivo
            slide_resumen = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content
            slide_resumen.shapes.title.text = "Resumen Ejecutivo"
            
            resumen = reporte.get('resumen_ejecutivo', {})
            content = slide_resumen.placeholders[1].text_frame
            content.text = f"Tipo de Contenido Ganador: {resumen.get('nombre_tipo', 'N/A')}"
            p = content.add_paragraph()
            p.text = f"Engagement Rate Promedio: {resumen.get('engagement_rate_promedio', 0):.2f}%"
            p = content.add_paragraph()
            p.text = f"Engagement Score Promedio: {resumen.get('engagement_score_promedio', 0):.1f}"
            if resumen.get('mejor_horario'):
                p = content.add_paragraph()
                p.text = f"Mejor Horario: {resumen['mejor_horario']}"
            if resumen.get('mejor_plataforma'):
                p = content.add_paragraph()
                p.text = f"Mejor Plataforma: {resumen['mejor_plataforma']}"
            
            # Slide 3: Insights Clave
            if 'analisis_ia' in reporte or hasattr(self.analizador_ai, 'analizar_con_ia'):
                analisis_ia = self.analizador_ai.analizar_con_ia(reporte)
                if 'analisis_ia' in analisis_ia:
                    slide_insights = prs.slides.add_slide(prs.slide_layouts[1])
                    slide_insights.shapes.title.text = "Insights Clave"
                    content = slide_insights.placeholders[1].text_frame
                    
                    insights = analisis_ia['analisis_ia'].get('insights_clave', [])
                    for i, insight in enumerate(insights[:5], 1):
                        if i == 1:
                            content.text = f"{i}. {insight}"
                        else:
                            p = content.add_paragraph()
                            p.text = f"{i}. {insight}"
            
            # Slide 4: Recomendaciones
            recomendaciones = self.analizador_ai.generar_recomendaciones_inteligentes(reporte)
            if recomendaciones:
                slide_recs = prs.slides.add_slide(prs.slide_layouts[1])
                slide_recs.shapes.title.text = "Recomendaciones Prioritarias"
                content = slide_recs.placeholders[1].text_frame
                
                for i, rec in enumerate(recomendaciones[:5], 1):
                    if i == 1:
                        content.text = f"[{rec['prioridad']}] {rec['titulo']}"
                    else:
                        p = content.add_paragraph()
                        p.text = f"[{rec['prioridad']}] {rec['titulo']}"
                    p = content.add_paragraph()
                    p.text = f"  {rec['accion']}"
                    p.level = 1
            
            # Slide 5: Métricas por Plataforma
            slide_plataformas = prs.slides.add_slide(prs.slide_layouts[1])
            slide_plataformas.shapes.title.text = "Métricas por Plataforma"
            content = slide_plataformas.placeholders[1].text_frame
            
            analisis_plataformas = reporte.get('analisis_por_plataforma', {})
            for plataforma, datos in list(analisis_plataformas.items())[:5]:
                if content.text == "":
                    content.text = f"{plataforma}:"
                else:
                    p = content.add_paragraph()
                    p.text = f"{plataforma}:"
                p = content.add_paragraph()
                p.text = f"  Engagement Rate: {datos.get('engagement_rate_promedio', 0):.2f}%"
                p.level = 1
            
            # Guardar presentación
            prs.save(output_file)
            
            return {
                "success": True,
                "archivo": output_file,
                "slides_generados": len(prs.slides),
                "timestamp": datetime.now().isoformat()
            }
            
        except Exception as e:
            logger.error(f"Error generando PowerPoint: {e}")
            return {"error": str(e)}
    
    def generar_dashboard_html(self, reporte: Dict[str, Any], output_file: str) -> Dict[str, Any]:
        """
        Genera un dashboard HTML interactivo con gráficos
        
        Args:
            reporte: Reporte generado
            output_file: Ruta del archivo HTML a generar
        
        Returns:
            Información del archivo generado
        """
        resumen = reporte.get('resumen_ejecutivo', {})
        
        html_content = f"""<!DOCTYPE html>
<html lang="es">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>Dashboard de Engagement</title>
    <script src="https://cdn.jsdelivr.net/npm/chart.js"></script>
    <style>
        * {{
            margin: 0;
            padding: 0;
            box-sizing: border-box;
        }}
        body {{
            font-family: 'Segoe UI', Tahoma, Geneva, Verdana, sans-serif;
            background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
            padding: 20px;
            min-height: 100vh;
        }}
        .container {{
            max-width: 1400px;
            margin: 0 auto;
        }}
        .header {{
            background: white;
            padding: 30px;
            border-radius: 10px;
            margin-bottom: 20px;
            box-shadow: 0 4px 6px rgba(0,0,0,0.1);
        }}
        .header h1 {{
            color: #333;
            margin-bottom: 10px;
        }}
        .header p {{
            color: #666;
        }}
        .grid {{
            display: grid;
            grid-template-columns: repeat(auto-fit, minmax(300px, 1fr));
            gap: 20px;
            margin-bottom: 20px;
        }}
        .card {{
            background: white;
            padding: 25px;
            border-radius: 10px;
            box-shadow: 0 4px 6px rgba(0,0,0,0.1);
        }}
        .card h2 {{
            color: #333;
            margin-bottom: 15px;
            font-size: 1.2em;
        }}
        .metric {{
            font-size: 2.5em;
            font-weight: bold;
            color: #667eea;
            margin-bottom: 10px;
        }}
        .metric-label {{
            color: #666;
            font-size: 0.9em;
        }}
        .chart-container {{
            position: relative;
            height: 300px;
            margin-top: 20px;
        }}
        .insights {{
            background: white;
            padding: 25px;
            border-radius: 10px;
            box-shadow: 0 4px 6px rgba(0,0,0,0.1);
            margin-bottom: 20px;
        }}
        .insights ul {{
            list-style: none;
            padding-left: 0;
        }}
        .insights li {{
            padding: 10px;
            margin-bottom: 10px;
            background: #f8f9fa;
            border-left: 4px solid #667eea;
            border-radius: 4px;
        }}
    </style>
</head>
<body>
    <div class="container">
        <div class="header">
            <h1>📊 Dashboard de Engagement</h1>
            <p>Reporte generado: {datetime.now().strftime('%d/%m/%Y %H:%M')}</p>
        </div>
        
        <div class="grid">
            <div class="card">
                <h2>Engagement Rate</h2>
                <div class="metric">{resumen.get('engagement_rate_promedio', 0):.2f}%</div>
                <div class="metric-label">Promedio</div>
            </div>
            <div class="card">
                <h2>Engagement Score</h2>
                <div class="metric">{resumen.get('engagement_score_promedio', 0):.1f}</div>
                <div class="metric-label">Puntos</div>
            </div>
            <div class="card">
                <h2>Mejor Plataforma</h2>
                <div class="metric">{resumen.get('mejor_plataforma', 'N/A')}</div>
                <div class="metric-label">Rendimiento óptimo</div>
            </div>
            <div class="card">
                <h2>Mejor Horario</h2>
                <div class="metric">{resumen.get('mejor_horario', 'N/A')}</div>
                <div class="metric-label">Timing óptimo</div>
            </div>
        </div>
        
        <div class="card">
            <h2>Engagement por Plataforma</h2>
            <div class="chart-container">
                <canvas id="platformChart"></canvas>
            </div>
        </div>
        
        <div class="insights">
            <h2>💡 Insights Clave</h2>
            <ul id="insights-list">
                <li>Tipo de contenido ganador: {resumen.get('nombre_tipo', 'N/A')}</li>
                <li>Total de publicaciones analizadas: {len(self.analizador.publicaciones)}</li>
            </ul>
        </div>
    </div>
    
    <script>
        // Gráfico de plataformas
        const platformData = {json.dumps(self._preparar_datos_grafico_plataformas(reporte), ensure_ascii=False)};
        
        const ctx = document.getElementById('platformChart').getContext('2d');
        new Chart(ctx, {{
            type: 'bar',
            data: {{
                labels: platformData.labels,
                datasets: [{{
                    label: 'Engagement Rate (%)',
                    data: platformData.data,
                    backgroundColor: 'rgba(102, 126, 234, 0.6)',
                    borderColor: 'rgba(102, 126, 234, 1)',
                    borderWidth: 2
                }}]
            }},
            options: {{
                responsive: true,
                maintainAspectRatio: false,
                scales: {{
                    y: {{
                        beginAtZero: true
                    }}
                }}
            }}
        }});
    </script>
</body>
</html>"""
        
        try:
            with open(output_file, 'w', encoding='utf-8') as f:
                f.write(html_content)
            
            return {
                "success": True,
                "archivo": output_file,
                "timestamp": datetime.now().isoformat()
            }
        except Exception as e:
            return {"error": str(e)}
    
    def _preparar_datos_grafico_plataformas(self, reporte: Dict[str, Any]) -> Dict[str, Any]:
        """Prepara datos para gráfico de plataformas"""
        analisis_plataformas = reporte.get('analisis_por_plataforma', {})
        
        labels = []
        data = []
        
        for plataforma, datos in analisis_plataformas.items():
            labels.append(plataforma)
            data.append(datos.get('engagement_rate_promedio', 0))
        
        return {
            "labels": labels,
            "data": data
        }
    
    def analizar_competencia(
        self,
        datos_competencia: List[Dict[str, Any]],
        metricas_propias: Dict[str, Any]
    ) -> Dict[str, Any]:
        """
        Analiza rendimiento vs competencia
        
        Args:
            datos_competencia: Lista de métricas de competidores
            metricas_propias: Métricas propias para comparar
        
        Returns:
            Análisis comparativo con competencia
        """
        if not datos_competencia:
            return {"error": "No hay datos de competencia"}
        
        # Calcular promedios de competencia
        engagement_rates_competencia = [c.get('engagement_rate', 0) for c in datos_competencia]
        engagement_scores_competencia = [c.get('engagement_score', 0) for c in datos_competencia]
        
        promedio_competencia = {
            'engagement_rate': statistics.mean(engagement_rates_competencia) if engagement_rates_competencia else 0,
            'engagement_score': statistics.mean(engagement_scores_competencia) if engagement_scores_competencia else 0
        }
        
        # Comparar con propias
        engagement_rate_propio = metricas_propias.get('engagement_rate', 0)
        engagement_score_propio = metricas_propias.get('engagement_score', 0)
        
        comparacion = {
            "metricas_propias": metricas_propias,
            "promedio_competencia": promedio_competencia,
            "posicion": {
                "engagement_rate": "por_encima" if engagement_rate_propio > promedio_competencia['engagement_rate'] else "por_debajo",
                "engagement_score": "por_encima" if engagement_score_propio > promedio_competencia['engagement_score'] else "por_debajo"
            },
            "diferencia": {
                "engagement_rate": engagement_rate_propio - promedio_competencia['engagement_rate'],
                "engagement_score": engagement_score_propio - promedio_competencia['engagement_score']
            },
            "percentil": self._calcular_percentil(
                engagement_rate_propio,
                engagement_rates_competencia
            )
        }
        
        return comparacion
    
    def _calcular_percentil(self, valor: float, valores_competencia: List[float]) -> int:
        """Calcula el percentil del valor propio vs competencia"""
        if not valores_competencia:
            return 50
        
        valores_ordenados = sorted(valores_competencia)
        posicion = sum(1 for v in valores_ordenados if v < valor)
        percentil = int((posicion / len(valores_ordenados)) * 100) if valores_ordenados else 50
        
        return percentil
    
    def generar_alertas_automaticas(self, reporte: Dict[str, Any]) -> List[Dict[str, Any]]:
        """
        Genera alertas automáticas basadas en umbrales
        
        Args:
            reporte: Reporte generado
        
        Returns:
            Lista de alertas generadas
        """
        alertas = []
        resumen = reporte.get('resumen_ejecutivo', {})
        
        # Alerta: Engagement rate bajo
        engagement_rate = resumen.get('engagement_rate_promedio', 0)
        if engagement_rate < 1.0:
            alertas.append({
                "nivel": "CRÍTICO",
                "tipo": "Engagement Rate Bajo",
                "mensaje": f"El engagement rate ({engagement_rate:.2f}%) está por debajo del umbral crítico (1.0%)",
                "accion": "Revisar estrategia de contenido inmediatamente",
                "timestamp": datetime.now().isoformat()
            })
        
        # Alerta: Tendencia decreciente
        if 'tendencia' in resumen and resumen['tendencia'].lower() == 'decreciente':
            alertas.append({
                "nivel": "ALTA",
                "tipo": "Tendencia Decreciente",
                "mensaje": "El engagement muestra tendencia decreciente",
                "accion": "Analizar causas y ajustar estrategia",
                "timestamp": datetime.now().isoformat()
            })
        
        # Alerta: Contenido viral bajo
        viral_pct = resumen.get('contenido_viral_porcentaje', 0)
        if viral_pct < 2.0:
            alertas.append({
                "nivel": "MEDIA",
                "tipo": "Bajo Contenido Viral",
                "mensaje": f"Solo {viral_pct:.1f}% del contenido alcanza criterios virales",
                "accion": "Optimizar contenido para mayor viralidad",
                "timestamp": datetime.now().isoformat()
            })
        
        return alertas
    
    def generar_reporte_ejecutivo(self, reporte: Dict[str, Any]) -> Dict[str, Any]:
        """
        Genera un reporte ejecutivo resumido para directivos
        
        Args:
            reporte: Reporte completo
        
        Returns:
            Reporte ejecutivo resumido
        """
        resumen = reporte.get('resumen_ejecutivo', {})
        
        # Obtener análisis con IA si está disponible
        insights_ia = []
        recomendaciones_ia = []
        
        try:
            analisis_ia = self.analizador_ai.analizar_con_ia(reporte)
            if 'analisis_ia' in analisis_ia:
                insights_ia = analisis_ia['analisis_ia'].get('insights_clave', [])[:3]
                recomendaciones_ia = analisis_ia['analisis_ia'].get('recomendaciones', [])[:3]
        except:
            pass
        
        reporte_ejecutivo = {
            "fecha": datetime.now().strftime('%Y-%m-%d'),
            "resumen": {
                "tipo_contenido_ganador": resumen.get('nombre_tipo', 'N/A'),
                "engagement_rate": resumen.get('engagement_rate_promedio', 0),
                "engagement_score": resumen.get('engagement_score_promedio', 0),
                "mejor_plataforma": resumen.get('mejor_plataforma', 'N/A'),
                "mejor_horario": resumen.get('mejor_horario', 'N/A')
            },
            "insights_clave": insights_ia or ["Análisis completo disponible en reporte detallado"],
            "recomendaciones_prioritarias": recomendaciones_ia or ["Revisar reporte completo para recomendaciones"],
            "alertas": self.generar_alertas_automaticas(reporte),
            "metricas_clave": {
                "total_publicaciones": len(self.analizador.publicaciones),
                "contenido_viral": resumen.get('contenido_viral_porcentaje', 0),
                "tendencia": resumen.get('tendencia', 'N/A')
            }
        }
        
        return reporte_ejecutivo


def main():
    """Función principal para demostración"""
    import argparse
    
    parser = argparse.ArgumentParser(description='Análisis de Engagement Avanzado')
    parser.add_argument('--publicaciones', type=int, default=30, help='Número de publicaciones')
    parser.add_argument('--powerpoint', type=str, help='Generar PowerPoint (ruta de salida)')
    parser.add_argument('--dashboard', type=str, help='Generar dashboard HTML (ruta de salida)')
    parser.add_argument('--reporte-ejecutivo', action='store_true', help='Generar reporte ejecutivo')
    
    args = parser.parse_args()
    
    # Crear analizadores
    analizador_base = AnalizadorEngagement()
    analizador_base.generar_datos_ejemplo(num_publicaciones=args.publicaciones)
    
    analizador_avanzado = AnalizadorEngagementAvanzado(analizador_base)
    
    # Generar reporte
    reporte = analizador_base.generar_reporte()
    
    # Exportar PowerPoint
    if args.powerpoint:
        print(f"\n📊 Generando PowerPoint: {args.powerpoint}")
        resultado = analizador_avanzado.exportar_powerpoint(reporte, args.powerpoint)
        if 'success' in resultado:
            print(f"✅ PowerPoint generado: {resultado['archivo']}")
            print(f"   Slides: {resultado['slides_generados']}")
        else:
            print(f"❌ Error: {resultado.get('error', 'Desconocido')}")
    
    # Generar Dashboard HTML
    if args.dashboard:
        print(f"\n📊 Generando Dashboard HTML: {args.dashboard}")
        resultado = analizador_avanzado.generar_dashboard_html(reporte, args.dashboard)
        if 'success' in resultado:
            print(f"✅ Dashboard generado: {resultado['archivo']}")
        else:
            print(f"❌ Error: {resultado.get('error', 'Desconocido')}")
    
    # Reporte ejecutivo
    if args.reporte_ejecutivo:
        print("\n📊 Generando Reporte Ejecutivo...")
        reporte_ejec = analizador_avanzado.generar_reporte_ejecutivo(reporte)
        print("\n" + "="*80)
        print("REPORTE EJECUTIVO")
        print("="*80)
        print(json.dumps(reporte_ejec, indent=2, ensure_ascii=False, default=str))


if __name__ == "__main__":
    main()


