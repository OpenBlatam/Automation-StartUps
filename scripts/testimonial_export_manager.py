#!/usr/bin/env python3
"""
Gestor de Exportación Avanzada para Testimonios
Exporta a múltiples formatos: PDF, Excel avanzado, PowerPoint, etc.
"""

import json
import logging
from typing import Dict, List, Any, Optional
from datetime import datetime
from pathlib import Path

logger = logging.getLogger(__name__)


class ExportManager:
    """Gestor de exportación avanzada"""
    
    def __init__(self):
        """Inicializa el gestor de exportación"""
        self.supported_formats = ['json', 'csv', 'txt', 'html', 'pdf', 'excel', 'pptx']
    
    def export_to_pdf(
        self,
        post_data: Dict[str, Any],
        output_file: str,
        include_charts: bool = True
    ) -> str:
        """
        Exporta a PDF usando reportlab
        
        Args:
            post_data: Datos del post
            output_file: Archivo de salida
            include_charts: Incluir gráficos (requiere matplotlib)
        
        Returns:
            Ruta del archivo generado
        """
        try:
            from reportlab.lib.pagesizes import letter, A4
            from reportlab.lib import colors
            from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
            from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle
            from reportlab.lib.units import inch
            
            output_path = Path(output_file)
            output_path.parent.mkdir(parents=True, exist_ok=True)
            
            doc = SimpleDocTemplate(str(output_path), pagesize=A4)
            story = []
            styles = getSampleStyleSheet()
            
            # Título
            title_style = ParagraphStyle(
                'CustomTitle',
                parent=styles['Heading1'],
                fontSize=24,
                textColor=colors.HexColor('#667eea'),
                spaceAfter=30
            )
            story.append(Paragraph("Análisis de Testimonio", title_style))
            story.append(Spacer(1, 0.2*inch))
            
            # Contenido del post
            story.append(Paragraph("Contenido Generado", styles['Heading2']))
            story.append(Paragraph(post_data.get('post_content', ''), styles['Normal']))
            story.append(Spacer(1, 0.2*inch))
            
            # Métricas principales
            story.append(Paragraph("Métricas Principales", styles['Heading2']))
            metrics_data = [
                ['Métrica', 'Valor'],
                ['Longitud', str(post_data.get('length', 0))],
                ['Hashtags', str(len(post_data.get('hashtags', [])))],
            ]
            
            if 'engagement_prediction' in post_data:
                pred = post_data['engagement_prediction']
                metrics_data.extend([
                    ['Score de Engagement', f"{pred.get('predicted_score', 0)}/100"],
                    ['Engagement Rate', f"{pred.get('predicted_engagement_rate', 0)}%"],
                    ['Confianza', f"{pred.get('confidence', 0)}"],
                ])
            
            if 'roi_calculation' in post_data:
                roi = post_data['roi_calculation']
                metrics_data.extend([
                    ['ROI', f"{roi.get('roi_percentage', 0)}%"],
                    ['Ingresos Estimados', f"${roi.get('estimated_revenue', 0):.2f}"],
                ])
            
            table = Table(metrics_data)
            table.setStyle(TableStyle([
                ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#667eea')),
                ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
                ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
                ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                ('FONTSIZE', (0, 0), (-1, 0), 14),
                ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
                ('BACKGROUND', (0, 1), (-1, -1), colors.beige),
                ('GRID', (0, 0), (-1, -1), 1, colors.black)
            ]))
            story.append(table)
            story.append(Spacer(1, 0.3*inch))
            
            # Hashtags
            if post_data.get('hashtags'):
                story.append(Paragraph("Hashtags", styles['Heading2']))
                hashtags_text = ', '.join(post_data['hashtags'])
                story.append(Paragraph(hashtags_text, styles['Normal']))
                story.append(Spacer(1, 0.2*inch))
            
            # Recomendaciones
            if 'engagement_prediction' in post_data:
                pred = post_data['engagement_prediction']
                if pred.get('recommendations'):
                    story.append(Paragraph("Recomendaciones", styles['Heading2']))
                    for rec in pred['recommendations'][:5]:
                        story.append(Paragraph(f"• {rec}", styles['Normal']))
            
            doc.build(story)
            logger.info(f"PDF exportado a: {output_file}")
            return str(output_path)
        except ImportError:
            logger.warning("reportlab no está instalado. Instala con: pip install reportlab")
            return self._export_to_text_fallback(post_data, output_file.replace('.pdf', '.txt'))
        except Exception as e:
            logger.error(f"Error al exportar PDF: {e}")
            return self._export_to_text_fallback(post_data, output_file.replace('.pdf', '.txt'))
    
    def export_to_excel_advanced(
        self,
        posts: List[Dict[str, Any]],
        output_file: str
    ) -> str:
        """
        Exporta a Excel avanzado con múltiples hojas
        
        Args:
            posts: Lista de posts a exportar
            output_file: Archivo de salida
        
        Returns:
            Ruta del archivo generado
        """
        try:
            import openpyxl
            from openpyxl.styles import Font, PatternFill, Alignment
            from openpyxl.chart import BarChart, Reference
            
            output_path = Path(output_file)
            output_path.parent.mkdir(parents=True, exist_ok=True)
            
            wb = openpyxl.Workbook()
            
            # Hoja 1: Resumen
            ws_summary = wb.active
            ws_summary.title = "Resumen"
            
            # Encabezados
            headers = ['Post ID', 'Plataforma', 'Longitud', 'Hashtags', 'Score', 'Engagement Rate', 'ROI %']
            ws_summary.append(headers)
            
            # Estilo de encabezados
            header_fill = PatternFill(start_color="667eea", end_color="667eea", fill_type="solid")
            header_font = Font(bold=True, color="FFFFFF")
            
            for cell in ws_summary[1]:
                cell.fill = header_fill
                cell.font = header_font
                cell.alignment = Alignment(horizontal="center")
            
            # Datos
            for i, post in enumerate(posts, start=2):
                row = [
                    f"Post_{i-1}",
                    post.get('platform', 'N/A'),
                    post.get('length', 0),
                    len(post.get('hashtags', [])),
                    post.get('engagement_prediction', {}).get('predicted_score', 0),
                    post.get('engagement_prediction', {}).get('predicted_engagement_rate', 0),
                    post.get('roi_calculation', {}).get('roi_percentage', 0)
                ]
                ws_summary.append(row)
            
            # Ajustar ancho de columnas
            for column in ws_summary.columns:
                max_length = 0
                column_letter = column[0].column_letter
                for cell in column:
                    try:
                        if len(str(cell.value)) > max_length:
                            max_length = len(str(cell.value))
                    except:
                        pass
                adjusted_width = min(max_length + 2, 50)
                ws_summary.column_dimensions[column_letter].width = adjusted_width
            
            # Hoja 2: Detalles
            ws_details = wb.create_sheet("Detalles")
            ws_details.append(['Campo', 'Valor'])
            
            for i, post in enumerate(posts):
                ws_details.append([f"=== POST {i+1} ===", ""])
                ws_details.append(['Contenido', post.get('post_content', '')])
                ws_details.append(['Hashtags', ', '.join(post.get('hashtags', []))])
                ws_details.append(['CTA', post.get('call_to_action', 'N/A')])
                ws_details.append(['', ''])  # Separador
            
            # Guardar
            wb.save(str(output_path))
            logger.info(f"Excel avanzado exportado a: {output_file}")
            return str(output_path)
        except ImportError:
            logger.warning("openpyxl no está instalado. Instala con: pip install openpyxl")
            return self.export_to_csv(posts, output_file.replace('.xlsx', '.csv'))
        except Exception as e:
            logger.error(f"Error al exportar Excel: {e}")
            return self.export_to_csv(posts, output_file.replace('.xlsx', '.csv'))
    
    def export_to_csv(
        self,
        posts: List[Dict[str, Any]],
        output_file: str
    ) -> str:
        """Exporta a CSV simple"""
        import csv
        
        output_path = Path(output_file)
        output_path.parent.mkdir(parents=True, exist_ok=True)
        
        with open(output_path, 'w', newline='', encoding='utf-8') as f:
            writer = csv.writer(f)
            writer.writerow(['Plataforma', 'Contenido', 'Hashtags', 'Score', 'Engagement Rate'])
            
            for post in posts:
                writer.writerow([
                    post.get('platform', ''),
                    post.get('post_content', ''),
                    ', '.join(post.get('hashtags', [])),
                    post.get('engagement_prediction', {}).get('predicted_score', 0),
                    post.get('engagement_prediction', {}).get('predicted_engagement_rate', 0)
                ])
        
        logger.info(f"CSV exportado a: {output_file}")
        return str(output_path)
    
    def export_to_powerpoint(
        self,
        post_data: Dict[str, Any],
        output_file: str
    ) -> str:
        """
        Exporta a PowerPoint
        
        Args:
            post_data: Datos del post
            output_file: Archivo de salida
        
        Returns:
            Ruta del archivo generado
        """
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            
            output_path = Path(output_file)
            output_path.parent.mkdir(parents=True, exist_ok=True)
            
            prs = Presentation()
            prs.slide_width = Inches(10)
            prs.slide_height = Inches(7.5)
            
            # Slide 1: Título
            slide1 = prs.slides.add_slide(prs.slide_layouts[0])
            title = slide1.shapes.title
            subtitle = slide1.placeholders[1]
            title.text = "Análisis de Testimonio"
            subtitle.text = f"Plataforma: {post_data.get('platform', 'N/A')}"
            
            # Slide 2: Contenido
            slide2 = prs.slides.add_slide(prs.slide_layouts[1])
            title2 = slide2.shapes.title
            title2.text = "Contenido Generado"
            content = slide2.placeholders[1]
            content.text = post_data.get('post_content', '')
            
            # Slide 3: Métricas
            slide3 = prs.slides.add_slide(prs.slide_layouts[1])
            title3 = slide3.shapes.title
            title3.text = "Métricas Principales"
            metrics = slide3.placeholders[1]
            
            metrics_text = f"Longitud: {post_data.get('length', 0)} caracteres\n"
            metrics_text += f"Hashtags: {len(post_data.get('hashtags', []))}\n"
            
            if 'engagement_prediction' in post_data:
                pred = post_data['engagement_prediction']
                metrics_text += f"Score: {pred.get('predicted_score', 0)}/100\n"
                metrics_text += f"Engagement Rate: {pred.get('predicted_engagement_rate', 0)}%\n"
            
            metrics.text = metrics_text
            
            prs.save(str(output_path))
            logger.info(f"PowerPoint exportado a: {output_file}")
            return str(output_path)
        except ImportError:
            logger.warning("python-pptx no está instalado. Instala con: pip install python-pptx")
            return self._export_to_text_fallback(post_data, output_file.replace('.pptx', '.txt'))
        except Exception as e:
            logger.error(f"Error al exportar PowerPoint: {e}")
            return self._export_to_text_fallback(post_data, output_file.replace('.pptx', '.txt'))
    
    def _export_to_text_fallback(
        self,
        post_data: Dict[str, Any],
        output_file: str
    ) -> str:
        """Exporta a texto como fallback"""
        output_path = Path(output_file)
        output_path.parent.mkdir(parents=True, exist_ok=True)
        
        with open(output_path, 'w', encoding='utf-8') as f:
            f.write("ANÁLISIS DE TESTIMONIO\n")
            f.write("=" * 60 + "\n\n")
            f.write(f"Plataforma: {post_data.get('platform', 'N/A')}\n")
            f.write(f"Longitud: {post_data.get('length', 0)} caracteres\n\n")
            f.write("CONTENIDO:\n")
            f.write("-" * 60 + "\n")
            f.write(post_data.get('post_content', '') + "\n\n")
            f.write("HASHTAGS:\n")
            f.write(", ".join(post_data.get('hashtags', [])) + "\n")
        
        return str(output_path)
    
    def export_all_formats(
        self,
        post_data: Dict[str, Any],
        base_filename: str,
        formats: Optional[List[str]] = None
    ) -> Dict[str, str]:
        """
        Exporta a todos los formatos disponibles
        
        Args:
            post_data: Datos del post
            base_filename: Nombre base del archivo
            formats: Formatos a exportar (default: todos)
        
        Returns:
            Dict con rutas de archivos generados
        """
        formats = formats or ['json', 'csv', 'txt', 'pdf', 'excel', 'pptx']
        exported_files = {}
        
        base_path = Path(base_filename)
        base_name = base_path.stem
        base_dir = base_path.parent
        
        for fmt in formats:
            try:
                if fmt == 'json':
                    output_file = base_dir / f"{base_name}.json"
                    with open(output_file, 'w', encoding='utf-8') as f:
                        json.dump(post_data, f, indent=2, ensure_ascii=False)
                    exported_files['json'] = str(output_file)
                
                elif fmt == 'csv':
                    output_file = base_dir / f"{base_name}.csv"
                    self.export_to_csv([post_data], str(output_file))
                    exported_files['csv'] = str(output_file)
                
                elif fmt == 'txt':
                    output_file = base_dir / f"{base_name}.txt"
                    self._export_to_text_fallback(post_data, str(output_file))
                    exported_files['txt'] = str(output_file)
                
                elif fmt == 'pdf':
                    output_file = base_dir / f"{base_name}.pdf"
                    self.export_to_pdf(post_data, str(output_file))
                    exported_files['pdf'] = str(output_file)
                
                elif fmt == 'excel':
                    output_file = base_dir / f"{base_name}.xlsx"
                    self.export_to_excel_advanced([post_data], str(output_file))
                    exported_files['excel'] = str(output_file)
                
                elif fmt == 'pptx':
                    output_file = base_dir / f"{base_name}.pptx"
                    self.export_to_powerpoint(post_data, str(output_file))
                    exported_files['pptx'] = str(output_file)
            except Exception as e:
                logger.warning(f"Error al exportar formato {fmt}: {e}")
        
        return exported_files


